# imports
import os
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from Summarization import summarise
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import fitz

# function for converting the pdf to powerpoint
def convert_pdf_to_pptx(input_pdf, output_pptx, powerpoint_title, powerpoint_sub_title, slide_theme):
    # create powerpoint
    powerpoint_pptx = Presentation()

    # opening pdf file to extract information from
    with pdfplumber.open(input_pdf) as pdf:
        pdf_document = fitz.open(input_pdf)

        # create title slide including title and subtitle
        slide = powerpoint_pptx.slides.add_slide(powerpoint_pptx.slide_layouts[0])
        left = top = Inches(0)
        pic = slide.shapes.add_picture(slide_theme, left, top, width=powerpoint_pptx.slide_width, height=powerpoint_pptx.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        title_box = slide.shapes.title
        title_box.text = powerpoint_title
        title_text_frame = title_box.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(40)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 120, 216)

        subtitle_box = slide.placeholders[1]
        subtitle_box.text = powerpoint_sub_title
        subtitle_text_frame = subtitle_box.text_frame
        subtitle_text_frame.paragraphs[0].font.size = Pt(18)
        subtitle_text_frame.paragraphs[0].font.bold = False
        subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(175, 123, 81)

        pdf_document = fitz.open(input_pdf)

        # loop through page by page in pdf file
        for page_number in range(len(pdf.pages)):

            # loading information from given page
            page = pdf_document.load_page(page_number)

            # extracting images from page
            image_list = page.get_images(full=True)

            # loop through every image on page
            for img_index in image_list:
                # extract images
                img = img_index[0]
                e_img = pdf_document.extract_image(img)
                extracted_image = e_img["image"]

                # create file path for image
                image_path = f"temp_img_{page_number}_{img}.png"
                with open(image_path, 'wb') as img_file:
                    img_file.write(extracted_image)

                # create slide for image
                image_slide = powerpoint_pptx.slides.add_slide(powerpoint_pptx.slide_layouts[5])

                left = top = Inches(0)
                image_slide_image = image_slide.shapes.add_picture(slide_theme, left, top, width=powerpoint_pptx.slide_width, height=powerpoint_pptx.slide_height)
                image_slide.shapes._spTree.remove(image_slide_image._element)
                image_slide.shapes._spTree.insert(2, image_slide_image._element)

                image_width = Inches(5)
                image_height = Inches(5)
                left = (powerpoint_pptx.slide_width - image_width) // 2
                top = (powerpoint_pptx.slide_height - image_height) // 2

                # import extracted image into slide
                image_slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

                # remove temp file path for image
                os.remove(image_path)

            # extracting text from page
            page = pdf.pages[page_number]
            page_text = page.extract_text()
            first_line, *remaining_lines = page_text.split('\n')

            # setting the title of slide and formatting it
            text_slide = powerpoint_pptx.slides.add_slide(powerpoint_pptx.slide_layouts[5])

            left = top = Inches(0)
            text_slide_image = text_slide.shapes.add_picture(slide_theme, left, top, width=powerpoint_pptx.slide_width, height=powerpoint_pptx.slide_height)
            text_slide.shapes._spTree.remove(text_slide_image._element)
            text_slide.shapes._spTree.insert(2, text_slide_image._element)

            title_box = text_slide.shapes.title
            title_box.text = first_line
            title_text_frame = title_box.text_frame
            title_text_frame.paragraphs[0].font.size = Pt(24)
            title_text_frame.paragraphs[0].font.bold = True
            title_text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 120, 216)

            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4)

            # creating a textbox on the slide for content
            textbox = text_slide.shapes.add_textbox(left, top, width, height)
            text_box = textbox.text_frame
            text_box.word_wrap = True

            # extracted content is given to api to have information summarized
            summarized_content = gpt_summarise('\n'.join(remaining_lines))

            # summarized information is formatted into a textbox on slide
            p = text_box.add_paragraph()
            p.text = summarized_content["text"]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(175, 123, 81)

    # saving and opening powerpoint when finished converting
    powerpoint_pptx.save(output_pptx)
